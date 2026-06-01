#!/usr/bin/env python3
"""
extract-text — Extract text content from any supported file

Usage:
  python3 scripts/extract-text.py <path> [options]

Options:
  --pages  <range>   PDF only: page range e.g. "1-3" or "2" (default: all)
  --sheet  <name>    XLSX only: sheet name (default: all sheets)
  --ocr              Use OCR for images (requires Tesseract)
  --lang   <lang>    OCR language (default: eng). e.g. tha, eng+tha

Supported: PDF, DOCX, XLSX, PPTX, images, TXT/MD/CSV

Output JSON:
  { path, type, text, pages?: [...] }
"""
import json
import os
import sys

sys.path.insert(0, os.path.dirname(__file__))
from lib.detect import detect

args = sys.argv[1:]
get = lambda flag: args[args.index(flag) + 1] if flag in args else ""
has = lambda flag: flag in args


def parse_page_range(spec, total):
    if not spec:
        return list(range(total))
    parts = spec.split("-")
    if len(parts) == 2:
        return list(range(int(parts[0]) - 1, int(parts[1])))
    return [int(parts[0]) - 1]


def extract_pdf(path):
    pages_arg = get("--pages")
    try:
        import pdfplumber
        pages_text = []
        with pdfplumber.open(path) as pdf:
            indices = parse_page_range(pages_arg, len(pdf.pages))
            for i in indices:
                if 0 <= i < len(pdf.pages):
                    text = pdf.pages[i].extract_text() or ""
                    pages_text.append({"page": i + 1, "text": text.strip()})
        full = "\n\n".join(p["text"] for p in pages_text if p["text"])
        return {"text": full, "pages": pages_text, "page_count": len(pages_text)}
    except ImportError:
        # Fallback to pypdf
        from pypdf import PdfReader
        r = PdfReader(path)
        indices = parse_page_range(pages_arg, len(r.pages))
        pages_text = []
        for i in indices:
            if 0 <= i < len(r.pages):
                text = r.pages[i].extract_text() or ""
                pages_text.append({"page": i + 1, "text": text.strip()})
        full = "\n\n".join(p["text"] for p in pages_text if p["text"])
        return {"text": full, "pages": pages_text, "page_count": len(pages_text)}


def extract_docx(path):
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Include table cells
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    return {"text": text, "paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}


def extract_xlsx(path):
    from openpyxl import load_workbook
    sheet_arg = get("--sheet")
    wb = load_workbook(path, read_only=True, data_only=True)
    results = {}
    names = [sheet_arg] if sheet_arg and sheet_arg in wb.sheetnames else wb.sheetnames
    for name in names:
        ws = wb[name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows_text.append("\t".join(cells))
        results[name] = "\n".join(rows_text)
    wb.close()
    full = "\n\n=== Sheet: ".join(
        f"{name} ===\n{text}" for name, text in results.items()
    )
    return {"text": full, "sheets": list(results.keys())}


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        slides_text.append({"slide": i + 1, "text": "\n".join(parts)})
    full = "\n\n".join(
        f"[Slide {s['slide']}]\n{s['text']}" for s in slides_text if s["text"]
    )
    return {"text": full, "slides": slides_text, "slide_count": len(slides_text)}


def _find_tesseract() -> str | None:
    """Auto-detect tesseract binary path cross-platform."""
    import platform
    import shutil

    system = platform.system()

    if system == "Windows":
        candidates = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            rf"C:\Users\{os.environ.get('USERNAME', '')}\AppData\Local\Programs\Tesseract-OCR\tesseract.exe",
        ]
        for p in candidates:
            if os.path.isfile(p):
                return p
    elif system == "Darwin":
        candidates = [
            "/opt/homebrew/bin/tesseract",   # Apple Silicon
            "/usr/local/bin/tesseract",       # Intel Mac
        ]
        for p in candidates:
            if os.path.isfile(p):
                return p
    # Linux + fallback: rely on PATH
    return shutil.which("tesseract")


def extract_image(path):
    if has("--ocr"):
        try:
            import pytesseract
            from PIL import Image

            # Auto-configure tesseract path (Windows needs this; macOS/Linux usually don't)
            tess_path = _find_tesseract()
            if tess_path:
                pytesseract.pytesseract.tesseract_cmd = tess_path
            else:
                return {
                    "error": (
                        "Tesseract not found. Install it:\n"
                        "  macOS:   brew install tesseract\n"
                        "  Windows: winget install UB-Mannheim.TesseractOCR\n"
                        "  Linux:   sudo apt install tesseract-ocr"
                    ),
                    "text": "",
                }

            import tempfile

            lang = get("--lang") or "eng"
            img = Image.open(path)

            # Normalize mode for best Tesseract compatibility
            if img.mode == "RGBA":
                bg = Image.new("RGB", img.size, (255, 255, 255))
                bg.paste(img, mask=img.split()[3])
                img = bg
            elif img.mode not in ("RGB", "L"):
                img = img.convert("RGB")

            # Save as PNG to a real (non-symlink) temp path.
            # On macOS /tmp → /private/tmp symlink; Leptonica can't follow symlinks,
            # so we resolve with os.path.realpath before creating the temp file.
            real_tmp = os.path.realpath(tempfile.gettempdir())
            tmp_fd, tmp_path = tempfile.mkstemp(suffix=".png", dir=real_tmp)
            os.close(tmp_fd)
            try:
                img.save(tmp_path, format="PNG")
                img.close()
                # Pass file path directly — pytesseract won't create another temp file
                text = pytesseract.image_to_string(tmp_path, lang=lang)
            finally:
                os.unlink(tmp_path)

            return {"text": text.strip(), "method": "ocr", "lang": lang,
                    "tesseract_path": tess_path}
        except ImportError:
            return {"error": "pytesseract not installed. Run: pip install pytesseract", "text": ""}
        except Exception as e:
            return {"error": str(e), "text": ""}
    else:
        from PIL import Image
        img = Image.open(path)
        result = {
            "text": "",
            "note": "Image file — use --ocr flag to extract text via OCR",
            "width": img.width,
            "height": img.height,
            "format": img.format or "",
        }
        img.close()
        return result


def extract_text_file(path):
    with open(path, encoding="utf-8-sig", errors="replace") as f:
        text = f.read()
    return {"text": text, "lines": len(text.splitlines())}


HANDLERS = {
    "pdf":   extract_pdf,
    "docx":  extract_docx,
    "doc":   extract_docx,
    "xlsx":  extract_xlsx,
    "xls":   extract_xlsx,
    "csv":   extract_text_file,
    "tsv":   extract_text_file,
    "pptx":  extract_pptx,
    "ppt":   extract_pptx,
    "image": extract_image,
    "text":  extract_text_file,
}


def main():
    file_args = [a for a in args if not a.startswith("--") and a != get("--pages")
                 and a != get("--sheet") and a != get("--lang")]
    if not file_args:
        print(json.dumps({"error": "Usage: extract-text.py <path> [--pages 1-3] [--ocr] [--lang eng]"}))
        sys.exit(1)

    path = file_args[0]
    if not os.path.isfile(path):
        print(json.dumps({"error": f"File not found: {path}"}))
        sys.exit(1)

    ftype = detect(path)
    handler = HANDLERS.get(ftype)
    if not handler:
        print(json.dumps({"error": f"Unsupported file type: {ftype} ({os.path.basename(path)})"}))
        sys.exit(1)

    try:
        result = handler(path)
        result["path"] = os.path.abspath(path)
        result["type"] = ftype
        print(json.dumps(result, ensure_ascii=False, indent=2))
    except Exception as e:
        print(json.dumps({"error": str(e), "path": path}))
        sys.exit(1)


if __name__ == "__main__":
    main()
